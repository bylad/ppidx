import os
import sys
import decimal
from bs4 import BeautifulSoup
import re
import requests
import xlrd
import docx
import win32com.client
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from tkinter import Tk
from tkinter import filedialog


monthe = ['январе', 'феврале', 'марте', 'апреле', 'мае', 'июне', 'июле',
          'августе', 'сентябре', 'октябре', 'ноябре', 'декабре']

months = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль',
          'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']


def dat(txt):
    """ Поиск по регулярному выражению.
        Например, из строки "О производстве в январе-мае 2019 года"
        будет выделена подстрока
            при ( |.) = "январе-мае 2019"
            при ( *)(.)( *) = "январе - мае 2019"
    """
    dt = re.compile("[яфмаисонд]([а-я]+[е])( *)(.)( *)[яфмаисонд]([а-я]+[е])( |.)\d{4}")
    # ( *)(.)( *)
    match = re.search(dt, txt)
    if match:
        return txt[match.start():match.end()]


def datastr(strlink):
    """ Вычленяем из даты числа для записи в виде 20191105
        dig0-год, dig1-месяц1, dig2-месяц2
    """
    dig = [int(s) for s in strlink.split() if s.isdigit()]
    for string in monthe:
        regex = re.compile(string)
        match = re.search(regex, strlink)
        if match:
            # global month
            month = strlink[match.start():match.end()]
            monthnum = '{:02}'.format(monthe.index(month) + 1)
            dig.append(int(monthnum))
    return dig


def news(newstext):
    """ Находим требуемую новость
    """
    # Добавляем headers, из-за ошибки при 1ом вызове alink = news()
    # AttributeError: 'NoneType' object has no attribute 'get'
    headers = {'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) \
               AppleWebKit/537.36 (KHTML, like Gecko) Chrome/62.0.3202.94 \
               Safari/537.36'}
    response = requests.get(url, headers=headers)
    soup = BeautifulSoup(response.text, 'html.parser')
    alink = soup.find('a', text=re.compile(newstext))
    if alink is None or alink == '':
        print(f'Новости о {newstext} нет. Программа завершена.')
        os.system("pause")
        sys.exit()
    return alink


def newsin(file_name, textin):
    """ В найденной новости скачиваем файл о промышленности НАО.
        При его наличии без закачки приступаем к обработке файла.
    """
    if os.path.exists(file_name):
        print(f'Файл {file_name} существует, скачивание пропущено.')
    else:
        prom = ''
        rlink = requests.get(a_link)
        soup = BeautifulSoup(rlink.text, 'html.parser')
        for link in soup.findAll('a', text=re.compile(textin)):
            prom = 'https://arhangelskstat.gks.ru' + link.get('href')
        file_ext = link.get('href')[-4:]  # выбираем последние 4 символа
        file_name = file_name + file_ext
        if file_ext == 'docx' or file_ext == 'DOCX':
            file_name = file_name + '.docx'
        if file_ext not in (['.doc', '.DOC']):
            print("Внимание! Проверьте на сайте тип скачиваемого файла.")
            print("Работа программы завершена.")
            os.system("pause")
            sys.exit()
            # Добавляем headers, т.к. получали ошибку
        # "TimeoutError: [WinError 10060]" при попытке скачать файл
        headers = {'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) \
                   AppleWebKit/537.36 (KHTML, like Gecko) Chrome/62.0.3202.94 \
                   Safari/537.36'}
        rfile = requests.get(prom, headers=headers)
        with open(file_name, 'wb') as f:
            f.write(rfile.content)
    return file_name


def doc2docx(basedir):
    """ Преобразуем файл DOC в DOCX
    """
    word = win32com.client.Dispatch("Word.application")
    for dir_path, dirs, files in os.walk(basedir):
        for file_name in files:
            file_path = os.path.join(dir_path, file_name)
            file_name, file_extension = os.path.splitext(file_path)
            if file_extension.lower() == '.doc':
                docx_file = '{0}{1}'.format(file_path, 'x')
                # Skip conversion where docx file already exists
                if not os.path.isfile(docx_file):
                    print('Преобразование в docx\n{0}\n'.format(file_path))
                    try:
                        word_doc = word.Documents.Open(file_path, False, False, False)
                        # Замена слеша в пути с / на \\, т.к. doc.SaveAs не отрабатывает /
                        docxf = re.sub('\/', '\\\\', docx_file)
                        word_doc.SaveAs2(docxf, FileFormat=16)
                        word_doc.Close()
                    except Exception:
                        print('Failed to Convert: {0}'.format(file_path))
    # word.Quit()


def prom_from_doc(doc, plist, marker):
    """ Создаем список с новыми значениями, из скачанного с Ростатата файла
    plist:  'Электроэнергия', 'Пар', 'Изделия хлебобулочные недлит', 'Молоко',
            'Нефть', 'Кондитерские', 'Масло', 'Оленина'
    marker: 0 - в новости период указан, как <январь-месяц год>
            1 - в новости период указан, как <январь год>
            2 - в новости период указан, как <год>
    """
    t = 1  # указываем 2-ую таблицу
    k = 1  # во 2-ой таблице указываем индекс столбца
    if plist[0] == 'Индекс':
        if marker != 2:
            t = 0  # указываем 1-ую таблицу
            k = 3  # для 1-ой таблицы меняем индекс столбца
        else:
            t = 0  # указываем 1-ую таблицу
    a = []
    table = doc.tables[t]
    if t == 0 and len(table.columns) != 4:
        print(f'Внимание! Изменилась структура таблицы, проверьте скачанный файл.')
        print(f'Количество столбцов = {len(table.columns)}, в январе должно быть 4.')
        os.system("pause")
    if t == 1 and len(table.columns) != 3:
        print(f'Внимание! Изменилась структура таблицы, проверьте скачанный файл.')
        print(f'Количество столбцов = {len(table.columns)}, в январе должно быть 4.')
        os.system("pause")
    for ri in range(len(table.rows)):
        for string in plist:
            regex = re.compile(string)
            match = re.search(regex, table.cell(ri, 0).text)
            if match:
                if table.cell(ri, k).text[:3] == "...":
                    fl_prom = float(0)
                else:
                    fl_prom = float(re.sub(',', '.', table.cell(ri, k).text))
                plist_index = plist.index(string)
                a.insert(plist_index, fl_prom)
    return a


def pptx_in():
    """ Указываем путь к презентации о промышленном производстве в НАО
    """
    parent = Tk()
    parent.withdraw()
    location = filedialog.askopenfilename(
        title="Выберите презентацию о промышленном производстве",
        filetypes=(("Powerpoint files", "*.pptx *.ppt"), ("all files", "*.*")))
    global prs
    prs = Presentation(location)
    directory = os.path.split(location)[0]
    # Поменять рабочий каталог на папку с файлом Powerpoint
    try:
        os.chdir(directory)
    except Exception:
        print("По указанному пути файл не может быть сохранен. \
              Скопируйте презентацию в доступную вам папку.")
        os.system("pause")
        sys.exit()
    return directory


def shape_upd(txt_frame, text, fontsize):
    txt_frame.clear()
    txt_frame.fit_text()
    p = txt_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(55, 96, 146)


def new_pptx(prsoutf):
    """ Заполнение презентации актуальными данными
        - zagolovok: замена указанных в заголовке месяца (года)
        - mes: замена указанных в надписи месяца (года)
        - index: замена индекса в %
        - prom0-7: замена соответствующих показателей
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            # Замена даты в заголовке
            if shape.name == 'zagolovok':
                text_frame = shape.text_frame
                cur_text = shape.text
                search_str = dat(shape.text)
                newdate = new_date
                # используем лок.перем.newdate во избежание ошибки присваивания внеш.new_date
                if len(newdate) < 5:  # только год, => меняем года на году
                    search_str = search_str + ' года'
                    newdate = newdate + ' году'
                new_text = cur_text.replace(search_str, newdate)
                shape_upd(text_frame, new_text, 26)
            # Замена даты в индексе
            if shape.name == 'mes':
                text_frame = shape.text_frame
                shape_upd(text_frame, new_mes, 16)
            if shape.name == 'index':
                text_frame = shape.text_frame
                indx = re.sub('\.', ',', str(idx[0])) + '%'
                shape_upd(text_frame, indx, 40)
            for k in range(8):
                if shape.name == 'prom' + str(k):
                    text_frame = shape.text_frame
                    if ap[k] == '0':
                        shape_upd(text_frame, '-', 40)
                    else:
                        shape_upd(text_frame, ap[k], 40)
    if os.path.exists(prsoutf):
        print('Файл с таким именем существует. Программа завершена.')
        os.system("pause")
        sys.exit()
    else:
        print(f'Сформирована презентация:\n{prsoutf}')
        prs.save(prsoutf)
        os.system("pause")


katalog = pptx_in()
url = 'https://arhangelskstat.gks.ru/news'

a_tag = news('промышленном производстве')
a_link = a_tag.get('href')
digs = datastr(a_tag.text)

try_marker = 0

try:
    digs[2]
except IndexError:
    try_marker = 2

try:
    digs[1]
except IndexError:
    try_marker = 1

if try_marker == 1:
    new_date = f'{str(digs[0])}'
    new_mes = f'{str(digs[0])} в %\n{str(digs[0] - 1)} г.'
    date_on_site = f'{digs[0]}_01-12'

if try_marker == 2:
    new_date = f'{monthe[digs[1] - 1]} {str(digs[0])}'
    new_mes = f'{months[digs[1] - 1]} {str(digs[0])} в %\n\
{months[digs[1] - 1]} {str(digs[0] - 1)} г.'
    date_on_site = f'{digs[0]}_{"{:02d}".format(digs[1])}-{"{:02d}".format(digs[1])}'

if try_marker == 0:
    new_date = f'{monthe[digs[1] - 1]}-{monthe[digs[2] - 1]} {str(digs[0])}'
    new_mes = f'{months[digs[1] - 1][:3]}-{months[digs[2] - 1]} {str(digs[0])} в %\n\
{months[digs[1] - 1][:3]}-{months[digs[2] - 1]} {str(digs[0] - 1)} г.'
    date_on_site = f'{digs[0]}_{"{:02d}".format(digs[1])}-{"{:02d}".format(digs[2])}'

path_year = f'{katalog}/{digs[0]}/'
if not os.path.exists(path_year):
    os.mkdir(path_year)

doc_file = newsin(path_year + f'{date_on_site}_prom', 'Ненецкому автономному округу')
if not os.path.exists(doc_file + 'x'):
    doc2docx(path_year)
docx_file = docx.Document(os.path.join(path_year, doc_file + 'x'))

'''
0-электроэнергия
1-пар
2-хлеб недлит.
3-молоко
4-нефть
5-кондитер
6-масло
7-оленина
'''
lprom = ['Электроэнергия', 'Пар', 'Изделия хлебобулочные недлит', 'Молоко',
         'Нефть', 'Кондитерские', 'Масло', 'Оленина']
a_values = prom_from_doc(docx_file, lprom, try_marker)
# try_marker-не важный параметр для a_values

index = ['Индекс']
idx = prom_from_doc(docx_file, index, try_marker)

i = 0
ap = [''] * 8
for i in range(4):
    ap[i] = str('{0:,}'.format(int(a_values[i] * 1000 // 1)).replace(',', ' '))
for i in range(4, 8):
    ap[i] = re.sub('\.', ',', str(a_values[i]))

prs_save = katalog + '/Stat_industry_' + date_on_site + '.pptx'
new_pptx(prs_save)
